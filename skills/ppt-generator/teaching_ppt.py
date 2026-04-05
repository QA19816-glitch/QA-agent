#!/usr/bin/env python3
"""
教学演示PPT生成器 - Skills技能清单教学演示
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# 颜色配置
COLORS = {
    'primary': RGBColor(147, 51, 234),
    'secondary': RGBColor(236, 72, 153),
    'accent': RGBColor(34, 211, 238),
    'highlight': RGBColor(250, 204, 21),
    'dark': RGBColor(30, 41, 59),
    'light': RGBColor(248, 250, 252),
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(51, 65, 85),
    'green': RGBColor(34, 197, 94),
    'orange': RGBColor(249, 115, 22),
}

def create_teaching_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ===== 第1页: 封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 装饰圆形
    for x, y, size, color in [(Inches(10), Inches(-1), Inches(4), COLORS['secondary']),
                               (Inches(-1), Inches(5), Inches(3), COLORS['accent'])]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "OpenClaw Skills 技能体系"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "完整技能清单与分类教学"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    info = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.5))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "测试专项 + 非测试专项 | 212个技能全面覆盖"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # ===== 第2页: 目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 课程大纲"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    content = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True
    
    chapters = [
        ("01", "资源入口概览", "GitHub官方地址与ZIP包下载"),
        ("02", "软件测试专项技能", "9大类93个测试技能详解"),
        ("03", "非测试专项技能", "7大类119个通用技能"),
        ("04", "数据统计分析", "技能分布与分类可视化"),
        ("05", "核心推荐清单", "建议安装的12个必备技能")
    ]
    
    for i, (num, ch_title, desc) in enumerate(chapters):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{num}  {ch_title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.space_before = Pt(16)
        
        p2 = tf.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['text']
    
    # ===== 第3页: 章节1 - 资源入口 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()
    
    num = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(3), Inches(2))
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "资源入口概览"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # ===== 第4页: GitHub资源入口 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "🐙 GitHub官方资源入口"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    content = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True
    
    items = [
        ("🏠 主页", "github.com/QA19816-glitch/QA-agent", "所有技能资源汇总"),
        ("🧪 测试专项", "SOFTWARE_TESTING_SPECIALTIES.md", "93个测试技能"),
        ("🌐 非测试专项", "NON_TESTING_SPECIALTIES.md", "119个通用技能"),
        ("✅ 建议安装", "ESSENTIAL_SETUP_SKILLS.md", "12个核心技能推荐"),
        ("📋 全局总清单", "ALL_SKILLS.md", "212个技能完整列表")
    ]
    
    for i, (icon, name, desc) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {name}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.space_before = Pt(14)
        
        p2 = tf.add_paragraph()
        p2.text = f"      {desc}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['text']
    
    # ===== 第5页: ZIP包下载 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light']
    bg.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "📦 分类ZIP包下载"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # 统计卡片
    cards = [
        ("93", "测试专项技能", COLORS['primary'], "🧪"),
        ("119", "非测试专项技能", COLORS['secondary'], "🌐"),
        ("12", "建议安装技能", COLORS['accent'], "✅")
    ]
    
    for i, (num, label, color, icon) in enumerate(cards):
        x = Inches(0.8 + i * 4.2)
        y = Inches(2)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.fill.background()
        
        icon_box = slide.shapes.add_textbox(x, y + Inches(0.3), Inches(3.8), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(40)