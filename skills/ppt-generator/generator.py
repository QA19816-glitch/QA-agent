#!/usr/bin/env python3
"""
PPT Generator Pro - 专业PPT生成器
支持图标、图表、数值展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# 颜色配置 - 活力多彩主题
COLORS = {
    'primary': RGBColor(147, 51, 234),
    'secondary': RGBColor(236, 72, 153),
    'accent': RGBColor(34, 211, 238),
    'highlight': RGBColor(250, 204, 21),
    'dark': RGBColor(30, 41, 59),
    'light': RGBColor(248, 250, 252),
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(51, 65, 85),
    'green': RGBColor(34, 197, 94),
    'orange': RGBColor(249, 115, 22),
}

class PPTGenerator:
    def __init__(self, template='vibrant'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.template = template
        
    def add_title_slide(self, title, subtitle, speaker="", date=""):
        """添加封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 紫色渐变背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['primary']
        bg.line.fill.background()
        
        # 装饰圆形
        for x, y, size, color in [
            (Inches(10), Inches(-1), Inches(4), COLORS['secondary']),
            (Inches(-1), Inches(5), Inches(3), COLORS['accent']),
        ]:
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS['accent']
            p.alignment = PP_ALIGN.CENTER
        
        # 演讲者信息
        if speaker or date:
            info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.6))
            tf = info_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{speaker}  |  {date}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['white']
            p.alignment = PP_ALIGN.CENTER
            
        return slide
    
    def add_section_slide(self, number, title, subtitle=""):
        """添加章节分隔页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 深色背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['dark']
        bg.line.fill.background()
        
        # 章节编号
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(3), Inches(2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{number}" if number < 10 else str(number)
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.6))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['accent']
        
        return slide
    
    def add_content_slide(self, title, bullets, accent_color='primary'):
        """添加内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 白色背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['white']
        bg.line.fill.background()
        
        # 左侧装饰条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), self.prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS[accent_color]
        bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        # 内容列表
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(10)
            p.line_spacing = 1.3
            
        return slide
    
    def add_stats_slide(self, title, stats):
        """添加数据统计页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 浅灰背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['light']
        bg.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        # 统计卡片
        card_width = Inches(3.8)
        card_height = Inches(2.2)
        gap = Inches(0.35)
        start_x = Inches(0.7)
        start_y = Inches(1.6)
        
        colors = [COLORS['primary'], COLORS['secondary'], COLORS['accent'], COLORS['highlight']]
        
        for i, (label, value, desc) in enumerate(stats):
            col = i % 3
            row = i // 3
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + Inches(0.4))
            
            # 卡片背景
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS['white']
            card.line.fill.background()
            
            # 数值
            val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_width - Inches(0.4), Inches(1))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = colors[i % len(colors)]
            
            # 标签
            label_box = slide.shapes.add_textbox(x + Inches(0.2),